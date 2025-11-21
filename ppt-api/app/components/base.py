from abc import ABC, abstractmethod
from typing import Any, Dict, Optional
from pptx.slide import Slide
from ..themes.theme import Theme, Style

class Component(ABC):
    """Base class for all slide components"""
    
    def __init__(self, theme: Optional[Theme] = None):
        self.theme = theme or Theme("default")

    @abstractmethod
    def render(self, slide: Slide, content: Dict[str, Any]) -> None:
        """
        Render the component on the slide with the given content
        
        Args:
            slide: The slide to render on
            content: Dictionary containing the content for the component
        """
        pass

    def _apply_style(self, text_frame, style_name: str) -> None:
        """Apply a named style from the theme to a text frame"""
        if self.theme:
            style = self.theme.get_style(style_name)
            if style:
                style.apply_to_text_frame(text_frame)

    def _get_safe_content(self, content: Dict[str, Any], key: str, default: Any = "") -> Any:
        """Safely get content from the content dictionary"""
        return content.get(key, default) 
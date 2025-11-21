from typing import Any, Dict, List
from pptx.util import Inches, Pt
from pptx.slide import Slide
from .base import Component
from pptx.enum.text import MSO_AUTO_SIZE, MSO_VERTICAL_ANCHOR
from pptx.dml.color import RGBColor

# Utility to clamp box within slide
SLIDE_WIDTH = 10
SLIDE_HEIGHT = 7.5

def clamp_box(left, top, width, height):
    left = max(0, min(left, SLIDE_WIDTH))
    top = max(0, min(top, SLIDE_HEIGHT))
    width = max(1, min(width, SLIDE_WIDTH - left))
    height = max(1, min(height, SLIDE_HEIGHT - top))
    return left, top, width, height

def configure_textbox_frame(text_frame):
    text_frame.word_wrap = True
    text_frame.auto_size = MSO_AUTO_SIZE.NONE
    text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
    text_frame.margin_top = Inches(0.05)
    text_frame.margin_bottom = Inches(0.05)
    text_frame.margin_left = Inches(0.05)
    text_frame.margin_right = Inches(0.05)

class HeaderWithImage(Component):
    """A slide layout with a header and an image below it"""
    
    def render(self, slide: Slide, content: Dict[str, Any]) -> None:
        # Add and style title
        title_text = self._get_safe_content(content, "title")
        if hasattr(slide.shapes, "title") and slide.shapes.title:
            title = slide.shapes.title
            title.text = title_text
            self._apply_style(title.text_frame, "title")
        else:
            title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
            title_box.text_frame.text = title_text
            self._apply_style(title_box.text_frame, "title")

        # Add image if provided
        image_path = self._get_safe_content(content, "image_path")
        if image_path:
            slide.shapes.add_picture(
                image_path,
                Inches(1),
                Inches(2),
                width=Inches(8),
                height=Inches(5)
            )

class BulletWithTitle(Component):
    """A slide layout with a title and bullet points"""
    
    def render(self, slide: Slide, content: Dict[str, Any]) -> None:
        # Add and style title
        title_text = self._get_safe_content(content, "title")
        if hasattr(slide.shapes, "title") and slide.shapes.title:
            title = slide.shapes.title
            title.text = title_text
            self._apply_style(title.text_frame, "title")
        else:
            title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
            title_box.text_frame.text = title_text
            self._apply_style(title_box.text_frame, "title")

        # Add bullet points
        points: List[str] = self._get_safe_content(content, "points", [])
        if points:
            tf = None
            try:
                bullet_shape = slide.shapes.placeholders[1]
                tf = bullet_shape.text_frame
                tf.clear()
            except (KeyError, IndexError, AttributeError):
                bullet_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(4))
                tf = bullet_box.text_frame
                tf.clear()
            for point in points:
                p = tf.add_paragraph()
                p.text = point
                p.level = 0
                self._apply_style(tf, "bullet")

class TwoColumnText(Component):
    """A slide layout with two columns of text"""

    def render(self, slide: Slide, content: Dict[str, Any]) -> None:
        # Title
        title_text = self._get_safe_content(content, "title")
        left = content.get("left", 1)
        top = content.get("top", 1)
        width = content.get("width", 4)
        height = content.get("height", 4)
        column_gap = content.get("column_gap", 0.5)
        left, top, width, height = clamp_box(left, top, width, height)

        if hasattr(slide.shapes, "title") and slide.shapes.title:
            title_shape = slide.shapes.title
            title_shape.text = title_text
            self._apply_style(title_shape.text_frame, "title")
        else:
            title_box = slide.shapes.add_textbox(Inches(left), Inches(top - 0.5), Inches(width * 2 + column_gap), Inches(1))
            title_box.text_frame.text = title_text
            self._apply_style(title_box.text_frame, "title")

        # Helper to configure textbox safely
        def configure_textbox(shape, text, style_name):
            tf = shape.text_frame
            tf.word_wrap = True
            tf.auto_size = MSO_AUTO_SIZE.NONE  # No auto resizing
            tf.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
            tf.margin_top = Inches(0.05)
            tf.margin_bottom = Inches(0.05)
            tf.margin_left = Inches(0.05)
            tf.margin_right = Inches(0.05)
            p = tf.paragraphs[0]
            p.text = text
            self._apply_style(tf, style_name)

        # Left column
        left_text = self._get_safe_content(content, "left_text")
        if left_text:
            left_box = slide.shapes.add_textbox(
                Inches(left),
                Inches(top + 1.0),  # Adjust top to leave space for title
                Inches(width),
                Inches(height)
            )
            configure_textbox(left_box, left_text, "body")

        # Right column
        right_text = self._get_safe_content(content, "right_text")
        if right_text:
            right_box = slide.shapes.add_textbox(
                Inches(left + width + column_gap),
                Inches(top + 1.0),
                Inches(width),
                Inches(height)
            )
            configure_textbox(right_box, right_text, "body")


class ComparisonTable(Component):
    """A table for comparing features, metrics, or options"""
    def render(self, slide: Slide, content: Dict[str, Any]) -> None:
        data: List[List[str]] = content.get("data", [])
        left = content.get("left", 1)
        top = content.get("top", 2)
        width = content.get("width", 8)
        height = content.get("height", 3)
        left, top, width, height = clamp_box(left, top, width, height)
        if data:
            rows = len(data)
            cols = len(data[0])
            table = slide.shapes.add_table(rows, cols, Inches(left), Inches(top), Inches(width), Inches(height)).table
            for r, row in enumerate(data):
                for c, cell in enumerate(row):
                    table.cell(r, c).text = str(cell)
                    # Optionally style header row
                    if r == 0:
                        for paragraph in table.cell(r, c).text_frame.paragraphs:
                            for run in paragraph.runs:
                                run.font.bold = True

class IconList(Component):
    """A list with icons (uses bullet points as icons for simplicity)"""
    def render(self, slide: Slide, content: Dict[str, Any]) -> None:
        items: List[str] = content.get("items", [])
        left = content.get("left", 1)
        top = content.get("top", 2)
        width = content.get("width", 8)
        height = content.get("height", 4)
        left, top, width, height = clamp_box(left, top, width, height)
        shape = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = shape.text_frame
        tf.clear()
        for item in items:
            p = tf.add_paragraph()
            p.text = item
            p.level = 0
            p.font.size = Pt(18)
            p.font.bold = False
            p.font.italic = False
            p.font.name = "Arial"
            p.bullet = True
            p.font.color.rgb = RGBColor(0, 0, 0)  # Example: black

class QuoteBlock(Component):
    """A stylized block for a quote, with optional attribution"""
    def render(self, slide: Slide, content: Dict[str, Any]) -> None:
        quote = content.get("quote", "")
        author = content.get("author", "")
        left = content.get("left", 1)
        top = content.get("top", 2)
        width = content.get("width", 8)
        height = content.get("height", 2)
        left, top, width, height = clamp_box(left, top, width, height)
        # Truncate quote if too long
        max_chars = int(width * height * 20)  # rough estimate
        if len(quote) > max_chars:
            quote = quote[:max_chars-3] + "..."
        shape = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = shape.text_frame
        configure_textbox_frame(tf)
        tf.clear()
        p = tf.add_paragraph()
        p.text = f'“{quote}”'
        p.font.size = Pt(24)
        p.font.italic = True
        if author:
            p = tf.add_paragraph()
            p.text = f'- {author}'
            p.font.size = Pt(16)
            p.font.italic = False
            p.font.bold = True

class Timeline(Component):
    """A horizontal timeline with labeled milestones"""
    def render(self, slide: Slide, content: Dict[str, Any]) -> None:
        milestones: List[str] = content.get("milestones", [])
        left = content.get("left", 1)
        top = content.get("top", 4)
        width = content.get("width", 8)
        height = content.get("height", 1)
        left, top, width, height = clamp_box(left, top, width, height)
        if not milestones:
            return
        step = width / max(1, len(milestones)-1)
        for i, label in enumerate(milestones):
            x = left + i * step
            shape = slide.shapes.add_textbox(Inches(x), Inches(top), Inches(1.2), Inches(height))
            tf = shape.text_frame
            tf.clear()
            p = tf.add_paragraph()
            p.text = label
            p.font.size = Pt(14)
            p.font.bold = True
            # Optionally, add a line/arrow between milestones (not implemented here)

class ProcessFlow(Component):
    """A series of labeled boxes/arrows to show a process or workflow"""
    def render(self, slide: Slide, content: Dict[str, Any]) -> None:
        steps: List[str] = content.get("steps", [])
        left = content.get("left", 1)
        top = content.get("top", 4)
        width = content.get("width", 8)
        height = content.get("height", 1)
        left, top, width, height = clamp_box(left, top, width, height)
        if not steps:
            return
        step_width = width / max(1, len(steps))
        for i, label in enumerate(steps):
            x = left + i * step_width
            shape = slide.shapes.add_textbox(Inches(x), Inches(top), Inches(step_width-0.2), Inches(height))
            tf = shape.text_frame
            tf.clear()
            p = tf.add_paragraph()
            p.text = label
            p.font.size = Pt(14)
            p.font.bold = True
            # Optionally, add arrows between boxes (not implemented here)

class StatisticHighlight(Component):
    """A large, bold number or percentage with a label and optional supporting text"""
    def render(self, slide: Slide, content: Dict[str, Any]) -> None:
        value = content.get("value", "")
        label = content.get("label", "")
        subtext = content.get("subtext", "")
        left = content.get("left", 3)
        top = content.get("top", 2)
        width = content.get("width", 4)
        height = content.get("height", 2)
        left, top, width, height = clamp_box(left, top, width, height)
        shape = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = shape.text_frame
        tf.clear()
        p = tf.add_paragraph()
        p.text = str(value)
        p.font.size = Pt(48)
        p.font.bold = True
        if label:
            p = tf.add_paragraph()
            p.text = label
            p.font.size = Pt(20)
            p.font.bold = False
        if subtext:
            p = tf.add_paragraph()
            p.text = subtext
            p.font.size = Pt(14)
            p.font.italic = True

class CalloutBox(Component):
    """A colored box with a key message, optionally with an icon"""
    def render(self, slide: Slide, content: Dict[str, Any]) -> None:
        message = content.get("message", "")
        left = content.get("left", 1)
        top = content.get("top", 6)
        width = content.get("width", 8)
        height = content.get("height", 1)
        left, top, width, height = clamp_box(left, top, width, height)
        shape = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        fill = shape.fill
        fill.solid()
        color = content.get("color", None)
        if isinstance(color, (list, tuple)) and len(color) == 3:
            fill.fore_color.rgb = RGBColor(*color)
        else:
            fill.fore_color.rgb = RGBColor(255, 215, 0)  # Default: gold
        tf = shape.text_frame
        configure_textbox_frame(tf)
        tf.clear()
        p = tf.add_paragraph()
        p.text = message
        p.font.size = Pt(20)
        p.font.bold = True

class SectionDivider(Component):
    """A slide with a big title and a visual separator"""
    def render(self, slide: Slide, content: Dict[str, Any]) -> None:
        title = content.get("title", "Section")
        left = content.get("left", 1)
        top = content.get("top", 3)
        width = content.get("width", 8)
        height = content.get("height", 2)
        left, top, width, height = clamp_box(left, top, width, height)
        shape = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = shape.text_frame
        tf.clear()
        p = tf.add_paragraph()
        p.text = title
        p.font.size = Pt(36)
        p.font.bold = True
        # Add a line below the title
        slide.shapes.add_shape(
            1,  # msoShapeLine
            Inches(left), Inches(top + 1.5), Inches(width), Inches(0.05)
        ) 
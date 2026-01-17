"""
API v1 routes - Deprecated endpoints
"""
from fastapi import APIRouter, File
from fastapi.responses import JSONResponse
from pydantic import BaseModel
from typing import List, Dict, Optional, Any
from pptx import Presentation
from pptx.util import Inches
import uuid
import io
import base64
from pptx.enum.shapes import MSO_SHAPE_TYPE

from ...themes.theme import THEMES
from ...components.layouts import (
    HeaderWithImage, BulletWithTitle, TwoColumnText,
    ComparisonTable, IconList, QuoteBlock, Timeline, ProcessFlow, StatisticHighlight, CalloutBox, SectionDivider
)
from ...utils.text_formatting import apply_markdown_to_text_frame, configure_textbox_frame

# This router will need access to global state from main.py
# We'll pass these as dependencies or use a state management approach
# For now, we'll import them - but they should be passed as dependencies in production

router = APIRouter(prefix="/api/v1", tags=["v1-deprecated"])

# Models
class ComponentContent(BaseModel):
    component_type: str
    content: Dict[str, Any]

class ThemeRequest(BaseModel):
    theme_name: str

class TextRequest(BaseModel):
    text: str
    left: Optional[float] = None
    top: Optional[float] = None
    width: Optional[float] = None
    height: Optional[float] = None

class BulletPointsRequest(BaseModel):
    points: List[str]
    left: Optional[float] = None
    top: Optional[float] = None
    width: Optional[float] = None
    height: Optional[float] = None

class TextBoxRequest(BaseModel):
    text: str
    left: float = 1.0
    top: float = 1.0
    width: float = 4.0
    height: float = 1.0

class SlideBase64Request(BaseModel):
    slideBase64: str

class TitlePositionRequest(BaseModel):
    left: float
    top: float
    width: float
    height: float


# Helper function to get global state (will be injected from main)
def get_state():
    """Get global state from main.py"""
    import sys
    import importlib
    # Import main module to access global state
    # This avoids circular import by importing at runtime
    main_module = sys.modules.get('app.main')
    if main_module is None:
        # If not already loaded, import it
        import app.main as main_module
    return main_module.prs, main_module.slide_map, main_module.current_theme


@router.post("/set/slideBase64")
def set_slide_base64(req: SlideBase64Request):
    """
    Initialize the presentation by clearing prs and loading from base64.
    This sets the current PowerPoint presentation we are working with.
    DEPRECATED - Use v2 API instead
    """
    import sys
    import app.main as main_module
    prs = main_module.prs
    slide_map = main_module.slide_map
    
    try:
        # Clear the current presentation and slide map
        prs = Presentation()
        slide_map.clear()
        
        # Decode the base64 string
        decoded_bytes = base64.b64decode(req.slideBase64)
        
        # Load the presentation from the decoded bytes
        buffer = io.BytesIO(decoded_bytes)
        prs = Presentation(buffer)
        
        # Optionally, populate slide_map with existing slides
        for i, slide in enumerate(prs.slides):
            slide_id = str(uuid.uuid4())
            slide_map[slide_id] = slide
        
        # Update global state
        main_module.prs = prs
        main_module.slide_map = slide_map
        
        return {
            "status": "ok",
            "message": "Presentation initialized from base64",
            "slide_count": len(prs.slides),
            "slide_ids": list(slide_map.keys())
        }
    except Exception as e:
        return {
            "status": "error",
            "error": str(e)
        }


@router.get("/themes")
def list_themes():
    """DEPRECATED - Use v2 API instead"""
    return {"themes": list(THEMES.keys())}


@router.post("/theme")
def set_theme(req: ThemeRequest):
    """DEPRECATED - Use v2 API instead"""
    if req.theme_name not in THEMES:
        return {"error": "theme not found"}
    import sys
    import app.main as main_module
    main_module.current_theme = THEMES[req.theme_name]
    return {"status": "ok"}


@router.post("/slide/{slide_id}/title")
def add_title(slide_id: str, req: TextRequest):
    """Add title to a slide - DEPRECATED - Use v2 API instead"""
    _, slide_map, current_theme = get_state()
    
    if slide_id not in slide_map:
        return {"error": "slide not found"}
    
    slide = slide_map[slide_id]
    if not slide.shapes.title:
        # If no title placeholder, create a text box
        if req.left is not None and req.top is not None:
            text_box = slide.shapes.add_textbox(
                Inches(req.left), Inches(req.top),
                Inches(req.width or 8), Inches(req.height or 1.5)
            )
            configure_textbox_frame(text_box.text_frame)
            apply_markdown_to_text_frame(text_box.text_frame, req.text)
            if current_theme:
                current_theme.get_style("title").apply_to_text_frame(text_box.text_frame)
        else:
            return {"error": "position required when no title placeholder exists"}
    else:
        configure_textbox_frame(slide.shapes.title.text_frame)
        apply_markdown_to_text_frame(slide.shapes.title.text_frame, req.text)
        if current_theme:
            current_theme.get_style("title").apply_to_text_frame(slide.shapes.title.text_frame)
    return {"status": "ok"}


@router.post("/slide/{slide_id}/subtitle")
def add_subtitle(slide_id: str, req: TextRequest):
    """Add subtitle to a slide - DEPRECATED - Use v2 API instead"""
    _, slide_map, current_theme = get_state()
    
    if slide_id not in slide_map:
        return {"error": "slide not found"}
    
    slide = slide_map[slide_id]
    # Find subtitle placeholder (usually index 1)
    subtitle = None
    for shape in slide.shapes:
        if hasattr(shape, "placeholder_format"):
            if shape.placeholder_format.idx == 1:
                subtitle = shape
                break
    
    if not subtitle:
        # If no subtitle placeholder, create a text box
        if req.left is not None and req.top is not None:
            text_box = slide.shapes.add_textbox(
                Inches(req.left), Inches(req.top),
                Inches(req.width or 8), Inches(req.height or 1)
            )
            configure_textbox_frame(text_box.text_frame)
            apply_markdown_to_text_frame(text_box.text_frame, req.text)
            if current_theme:
                current_theme.get_style("subtitle").apply_to_text_frame(text_box.text_frame)
        else:
            return {"error": "position required when no subtitle placeholder exists"}
    else:
        configure_textbox_frame(subtitle.text_frame)
        apply_markdown_to_text_frame(subtitle.text_frame, req.text)
        if current_theme:
            current_theme.get_style("subtitle").apply_to_text_frame(subtitle.text_frame)
    return {"status": "ok"}


@router.post("/slide/{slide_id}/bullet_points")
def add_bullet_points(slide_id: str, req: BulletPointsRequest):
    """Add bullet points to a slide - DEPRECATED - Use v2 API instead"""
    _, slide_map, current_theme = get_state()
    
    if slide_id not in slide_map:
        return {"error": "slide not found"}
    
    slide = slide_map[slide_id]
    
    left = req.left if req.left is not None else 1
    top = req.top if req.top is not None else 2
    width = req.width if req.width is not None else 8
    height = req.height if req.height is not None else 4

    shape = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = shape.text_frame
    configure_textbox_frame(tf)
    
    markdown_text = "\n".join(req.points)
    apply_markdown_to_text_frame(tf, markdown_text)

    if current_theme:
        current_theme.get_style("bullet").apply_to_text_frame(tf)

    return {"status": "ok"}


@router.post("/slide/{slide_id}/text_box")
def add_text_box(slide_id: str, req: TextBoxRequest):
    """Add a text box to a slide - DEPRECATED - Use v2 API instead"""
    _, slide_map, current_theme = get_state()
    
    if slide_id not in slide_map:
        return {"error": "slide not found"}
    
    slide = slide_map[slide_id]
    left = req.left if req.left is not None else 1
    top = req.top if req.top is not None else 1
    width = req.width if req.width is not None else 8
    height = req.height if req.height is not None else 1
    text_box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    configure_textbox_frame(text_box.text_frame)
    apply_markdown_to_text_frame(text_box.text_frame, req.text)
    if current_theme:
        current_theme.get_style("body").apply_to_text_frame(text_box.text_frame)
    return {"status": "ok"}


@router.post("/slide/{slide_id}/component")
def add_component(slide_id: str, req: ComponentContent):
    """DEPRECATED - Use v2 API instead"""
    _, slide_map, current_theme = get_state()
    
    if slide_id not in slide_map:
        return {"error": "slide not found"}
    
    slide = slide_map[slide_id]
    component_map = {
        "header_with_image": HeaderWithImage,
        "bullet_with_title": BulletWithTitle,
        "two_column_text": TwoColumnText,
        "comparison_table": ComparisonTable,
        "icon_list": IconList,
        "quote_block": QuoteBlock,
        "timeline": Timeline,
        "process_flow": ProcessFlow,
        "statistic_highlight": StatisticHighlight,
        "callout_box": CalloutBox,
        "section_divider": SectionDivider
    }
    
    if req.component_type not in component_map:
        return {"error": "invalid component type"}
    
    component_class = component_map[req.component_type]
    component = component_class(theme=current_theme)
    component.render(slide, req.content)
    
    return {"status": "ok"}


@router.post("/slide")
def create_slide(layout: int = 0):
    """Create a new blank slide - DEPRECATED - Use v2 API instead"""
    prs, slide_map, current_theme = get_state()
    
    slide = prs.slides.add_slide(prs.slide_layouts[layout])
    slide_id = str(uuid.uuid4())
    slide_map[slide_id] = slide
    
    # Apply current theme to new slide
    current_theme.apply_to_slide(slide)
    
    return {"slide_id": slide_id}


@router.post("/slide/blank")
def create_blank_slide():
    """Create a new blank slide without any placeholders - DEPRECATED - Use v2 API instead"""
    prs, slide_map, current_theme = get_state()
    
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Layout 6 is blank
    slide_id = str(uuid.uuid4())
    slide_map[slide_id] = slide
    
    # Apply current theme to new slide
    current_theme.apply_to_slide(slide)
    
    return {"slide_id": slide_id}


@router.delete("/slide/{slide_id}")
def delete_slide(slide_id: str):
    """Delete a slide - DEPRECATED - Use v2 API instead"""
    prs, slide_map, _ = get_state()
    
    if slide_id in slide_map:
        slide = slide_map.pop(slide_id)
        # Find the slide in the presentation and remove it
        for idx, sld in enumerate(prs.slides):
            if sld == slide:
                xml_slides = prs.slides._sldIdLst
                xml_slides.remove(xml_slides[idx])
                return {"status": "deleted"}
        return {"error": "slide not found in presentation"}
    return {"error": "slide not found"}


@router.get("/slides")
def list_slides():
    """List all slides - DEPRECATED - Use v2 API instead"""
    _, slide_map, _ = get_state()
    return {"slide_ids": list(slide_map.keys())}


@router.post("/save")
def save_presentation(filename: str = "output.pptx"):
    """Save the presentation - DEPRECATED - Use v2 API instead"""
    prs, _, _ = get_state()
    prs.save(filename)
    return {"status": "saved", "filename": filename}


@router.get("/presentation/base64")
def get_presentation_base64():
    """Export the current presentation as base64 string - DEPRECATED - Use v2 API instead"""
    prs, _, _ = get_state()
    
    try:
        buffer = io.BytesIO()
        prs.save(buffer)
        buffer.seek(0)
        
        base64_data = base64.b64encode(buffer.read()).decode('utf-8')
        
        return {
            "status": "ok",
            "base64": base64_data,
            "filename": "presentation.pptx",
            "slide_count": len(prs.slides)
        }
    except Exception as e:
        return JSONResponse(
            status_code=500,
            content={"error": str(e)}
        )


@router.get("/presentation/preview")
def get_presentation_preview():
    """Get presentation metadata and base64 for preview - DEPRECATED - Use v2 API instead"""
    prs, slide_map, current_theme = get_state()
    
    try:
        buffer = io.BytesIO()
        prs.save(buffer)
        buffer.seek(0)
        
        base64_data = base64.b64encode(buffer.read()).decode('utf-8')
        
        return {
            "status": "ok",
            "base64": base64_data,
            "slide_count": len(prs.slides),
            "slide_ids": list(slide_map.keys()),
            "current_theme": current_theme.name if hasattr(current_theme, 'name') else "default"
        }
    except Exception as e:
        return JSONResponse(
            status_code=500,
            content={"error": str(e)}
        )


@router.post("/presentation/reset")
def reset_presentation():
    """Reset the presentation to start fresh - DEPRECATED - Use v2 API instead"""
    import sys
    import app.main as main_module
    main_module.prs = Presentation()
    main_module.slide_map = {}
    return {"status": "ok", "message": "Presentation reset"}


@router.get("/slide/{slide_id}/shapes")
def get_slide_shapes(slide_id: str):
    """Get detailed information about all shapes in a slide - DEPRECATED - Use v2 API instead"""
    _, slide_map, _ = get_state()
    
    if slide_id not in slide_map:
        return {"error": "slide not found"}
    
    slide = slide_map[slide_id]
    shapes_info = []
    
    for idx, shape in enumerate(slide.shapes):
        shape_info = {
            "shape_id": f"shape_{idx}",
            "shape_index": idx,
            "shape_name": shape.name,
            "has_text": shape.has_text_frame,
            "text_content": "",
            "position": {
                "left": shape.left.inches,
                "top": shape.top.inches,
                "width": shape.width.inches,
                "height": shape.height.inches
            },
            "is_title": False,
            "is_subtitle": False,
            "shape_type": "unknown"
        }
        
        # Determine shape type
        if shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
            shape_info["shape_type"] = "placeholder"
            if hasattr(shape, "placeholder_format"):
                if shape.placeholder_format.type == 1:  # Title
                    shape_info["is_title"] = True
                    shape_info["shape_type"] = "title"
                elif shape.placeholder_format.type == 2:  # Subtitle/Body
                    shape_info["is_subtitle"] = True
                    shape_info["shape_type"] = "subtitle"
        elif shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
            shape_info["shape_type"] = "text_box"
        elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            shape_info["shape_type"] = "picture"
        elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            shape_info["shape_type"] = "table"
        elif shape.shape_type == MSO_SHAPE_TYPE.CHART:
            shape_info["shape_type"] = "chart"
        elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            shape_info["shape_type"] = "group"
        elif shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            shape_info["shape_type"] = "auto_shape"
        
        if hasattr(slide.shapes, "title") and slide.shapes.title == shape:
            shape_info["is_title"] = True
            shape_info["shape_type"] = "title"
        
        if shape.has_text_frame:
            text_parts = []
            for paragraph in shape.text_frame.paragraphs:
                text_parts.append(paragraph.text)
            shape_info["text_content"] = "\n".join(text_parts)
        
        shapes_info.append(shape_info)
    
    return {
        "slide_id": slide_id,
        "shape_count": len(shapes_info),
        "shapes": shapes_info
    }


@router.post("/slide/{slide_id}/title/position")
def set_title_position(slide_id: str, req: TitlePositionRequest):
    """Set the position of the title shape on a slide - DEPRECATED - Use v2 API instead"""
    _, slide_map, _ = get_state()
    
    if slide_id not in slide_map:
        return {"error": "slide not found"}
    
    slide = slide_map[slide_id]
    
    # Find the title shape
    title_shape = None
    
    if hasattr(slide.shapes, "title") and slide.shapes.title:
        title_shape = slide.shapes.title
    else:
        for shape in slide.shapes:
            if hasattr(shape, "placeholder_format"):
                if shape.placeholder_format.type == 1:  # Title placeholder
                    title_shape = shape
                    break
    
    if not title_shape:
        return {
            "error": "title shape not found on this slide",
            "message": "This slide may not have a title placeholder"
        }
    
    # Update position
    title_shape.left = Inches(req.left)
    title_shape.top = Inches(req.top)
    title_shape.width = Inches(req.width)
    title_shape.height = Inches(req.height)
    
    return {
        "status": "ok",
        "message": "Title position updated successfully",
        "slide_id": slide_id,
        "new_position": {
            "left": req.left,
            "top": req.top,
            "width": req.width,
            "height": req.height
        }
    }


@router.get("/slide/{slide_id}/title/coordinates")
def get_title_coordinates(slide_id: str):
    """Get the position and dimensions of the title shape - DEPRECATED - Use v2 API instead"""
    _, slide_map, _ = get_state()
    
    from ...utils.shape_alignment import get_title_coordinates as get_title_coords_util
    return get_title_coords_util(slide_map, slide_id)


@router.post("/slides/bulk/title/position")
def set_bulk_title_positions(request: dict):
    """Set the same title position for multiple slides - DEPRECATED - Use v2 API instead"""
    _, slide_map, _ = get_state()
    
    from ...utils.shape_alignment import set_bulk_title_positions as set_bulk_titles_util
    return set_bulk_titles_util(slide_map, request)


@router.post("/slides/align_titles_to_reference")
def align_titles_to_reference(request: dict):
    """Align title positions to match a reference slide - DEPRECATED - Use v2 API instead"""
    _, slide_map, _ = get_state()
    
    from ...utils.shape_alignment import align_titles_to_reference as align_titles_util
    return align_titles_util(slide_map, request["reference_slide_number"], request["target_slide_numbers"])


@router.post("/slides/align_subtitles_to_reference")
def align_subtitles_to_reference(request: dict):
    """Align subtitle positions to match a reference slide - DEPRECATED - Use v2 API instead"""
    _, slide_map, _ = get_state()
    
    from ...utils.shape_alignment import align_subtitles_to_reference as align_subtitles_util
    return align_subtitles_util(slide_map, request["reference_slide_number"], request["target_slide_numbers"])


@router.post("/slides/align_footnotes_to_reference")
def align_footnotes_to_reference(request: dict):
    """Align footnote positions to match a reference slide - DEPRECATED - Use v2 API instead"""
    _, slide_map, _ = get_state()
    
    from ...utils.shape_alignment import align_footnotes_to_reference as align_footnotes_util
    return align_footnotes_util(slide_map, request["reference_slide_number"], request["target_slide_numbers"])


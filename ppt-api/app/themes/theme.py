from typing import Dict, Optional
from pptx.util import Pt
from pptx.text.text import _Paragraph
from pptx.dml.color import RGBColor

class Style:
    def __init__(
        self,
        font_size: float = 18,
        font_family: str = "Calibri",
        bold: bool = False,
        italic: bool = False,
        color: str = "#000000"
    ):
        self.font_size = font_size
        self.font_family = font_family
        self.bold = bold
        self.italic = italic
        self.color = color

    def apply_to_text_frame(self, text_frame) -> None:
        """Apply this style to a text frame"""
        for paragraph in text_frame.paragraphs:
            self.apply_to_paragraph(paragraph)

    def apply_to_paragraph(self, paragraph: _Paragraph) -> None:
        """Apply this style to a paragraph"""
        font = paragraph.font
        font.size = Pt(self.font_size)
        font.name = self.font_family
        font.bold = self.bold
        font.italic = self.italic
        
        # Convert hex color to RGB
        rgb = tuple(int(self.color.lstrip('#')[i:i+2], 16) for i in (0, 2, 4))
        font.color.rgb = RGBColor(*rgb)

class Theme:
    def __init__(
        self,
        name: str,
        font_family: str = "Calibri",
        primary_color: str = "#1F497D",
        background_color: str = "#FFFFFF"
    ):
        self.name = name
        self.font_family = font_family
        self.primary_color = primary_color
        self.background_color = background_color
        
        # Initialize default styles
        self.styles: Dict[str, Style] = {
            "title": Style(
                font_size=44,
                font_family=font_family,
                bold=True,
                color=primary_color
            ),
            "subtitle": Style(
                font_size=32,
                font_family=font_family,
                color=primary_color
            ),
            "heading": Style(
                font_size=28,
                font_family=font_family,
                bold=True,
                color=primary_color
            ),
            "body": Style(
                font_size=18,
                font_family=font_family,
                color="#000000"
            ),
            "bullet": Style(
                font_size=18,
                font_family=font_family,
                color="#000000"
            )
        }

    def get_style(self, style_name: str) -> Optional[Style]:
        """Get a style by name"""
        return self.styles.get(style_name)

    def apply_to_slide(self, slide) -> None:
        """Apply theme to a slide"""
        # Set background color
        background = slide.background
        fill = background.fill
        fill.solid()
        rgb = tuple(int(self.background_color.lstrip('#')[i:i+2], 16) for i in (0, 2, 4))
        fill.fore_color.rgb = RGBColor(*rgb)

        # Apply title style if title exists
        if hasattr(slide.shapes, "title") and slide.shapes.title:
            self.styles["title"].apply_to_text_frame(slide.shapes.title.text_frame)

# Create some predefined themes
THEMES = {
    "default": Theme(
        name="Default",
        primary_color="#1F497D",
        background_color="#FFFFFF"
    ),
    "dark": Theme(
        name="Dark",
        primary_color="#FFFFFF",
        background_color="#2F2F2F"
    ),
    "modern": Theme(
        name="Modern",
        font_family="Segoe UI",
        primary_color="#0078D4",
        background_color="#F5F5F5"
    ),
    # 🟦 Theme from CFI Investment Banking Pitchbook
    "pitchbook": Theme(
        name="Pitchbook",
        font_family="Arial",
        primary_color="#0B2341",       # Deep navy blue (common in finance decks)
        background_color="#FFFFFF"
    ),
    # 🟩 Theme from Placeholder Strategy Deck (PDF-based)
    "strategy_template": Theme(
        name="Strategy Template",
        font_family="Calibri",
        primary_color="#2F5597",       # Microsoft blue tone
        background_color="#FFFFFF"
    )
} 

THEMES["pitchbook"].styles["title"].font_size = 40
THEMES["pitchbook"].styles["subtitle"].font_size = 28
THEMES["pitchbook"].styles["heading"].font_size = 24
THEMES["pitchbook"].styles["body"].font_size = 16
THEMES["pitchbook"].styles["bullet"].font_size = 16

THEMES["strategy_template"].styles["title"].font_size = 36
THEMES["strategy_template"].styles["subtitle"].font_size = 24
THEMES["strategy_template"].styles["heading"].font_size = 22
THEMES["strategy_template"].styles["body"].font_size = 16
THEMES["strategy_template"].styles["bullet"].font_size = 16

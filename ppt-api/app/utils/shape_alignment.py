"""
Utility functions for shape alignment operations.
These are helper functions used by the alignment APIs but not exposed as APIs themselves.
"""
from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.shapes import PP_PLACEHOLDER
from typing import Dict, Any


def _find_subtitle(shapes):
    """
    Find a subtitle shape using BOTH:
    - Official python-pptx placeholder type (PP_PLACEHOLDER.SUBTITLE)
    - User-defined working logic: placeholder_format.type == 2
    """
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
            if hasattr(shape, "placeholder_format"):
                ph_type = shape.placeholder_format.type

                # Official subtitle placeholder
                if ph_type == PP_PLACEHOLDER.SUBTITLE:
                    return shape

                # User custom logic (your working system)
                # type == 2 means Subtitle/Body for your template
                if ph_type == 2:
                    return shape
    return None


def align_titles_to_reference(slide_map: Dict[str, Any], reference_slide_num: int, target_slide_numbers: list):
    """
    Align title positions of multiple slides to match a reference slide.
    
    This is a utility function, not an API endpoint.
    """
    # Get slide_ids_by_index
    slide_ids_by_index = list(slide_map.keys())
    total_slides = len(slide_ids_by_index)
    
    # Validate reference slide number
    ref_index = reference_slide_num - 1
    if ref_index < 0 or ref_index >= total_slides:
        return {
            "error": f"Invalid reference slide number. Valid range is 1-{total_slides}"
        }
    
    # Step 1: Get reference slide coordinates
    ref_slide_id = slide_ids_by_index[ref_index]
    ref_coordinates_response = get_title_coordinates(slide_map, ref_slide_id)
    
    # Check if reference slide has a title
    if "error" in ref_coordinates_response or not ref_coordinates_response.get("has_title"):
        return {
            "error": f"Reference slide {reference_slide_num} does not have a title shape",
            "details": ref_coordinates_response
        }
    
    # Extract position from reference
    reference_position = ref_coordinates_response["coordinates"]
    
    # Step 2: Apply position to target slides
    bulk_request = {
        "slide_numbers": target_slide_numbers,
        "position": reference_position
    }
    
    bulk_response = set_bulk_title_positions(slide_map, bulk_request)
    
    return {"data": "Success"}


def align_subtitles_to_reference(slide_map: Dict[str, Any], reference_slide_num: int, target_slide_numbers: list):
    """
    Align subtitle positions of multiple slides to match a reference slide.
    
    This is a utility function, not an API endpoint.
    """
    # Slide index handling
    slide_ids_by_index = list(slide_map.keys())
    total_slides = len(slide_ids_by_index)

    ref_index = reference_slide_num - 1
    if ref_index < 0 or ref_index >= total_slides:
        return {
            "error": f"Invalid reference slide number. Valid range is 1-{total_slides}"
        }

    # Get the reference slide object
    ref_slide = slide_map[slide_ids_by_index[ref_index]]

    # Find reference subtitle
    ref_subtitle_shape = _find_subtitle(ref_slide.shapes)

    if ref_subtitle_shape is None:
        return {
            "error": f"Reference slide {reference_slide_num} does not contain a subtitle placeholder"
        }

    # Extract reference subtitle geometry
    reference_position = {
        "left": ref_subtitle_shape.left.inches,
        "top": ref_subtitle_shape.top.inches,
        "width": ref_subtitle_shape.width.inches,
        "height": ref_subtitle_shape.height.inches
    }

    # Apply to targets
    results = []

    for slide_num in target_slide_numbers:
        slide_index = slide_num - 1

        if slide_index < 0 or slide_index >= total_slides:
            results.append({
                "slide_number": slide_num,
                "success": False,
                "error": f"Invalid slide number. Valid range is 1-{total_slides}"
            })
            continue

        slide = slide_map[slide_ids_by_index[slide_index]]

        subtitle_shape = _find_subtitle(slide.shapes)

        if subtitle_shape is None:
            results.append({
                "slide_number": slide_num,
                "success": False,
                "error": "No subtitle placeholder found on this slide"
            })
            continue

        try:
            subtitle_shape.left = Inches(reference_position["left"])
            subtitle_shape.top = Inches(reference_position["top"])
            subtitle_shape.width = Inches(reference_position["width"])
            subtitle_shape.height = Inches(reference_position["height"])

            results.append({
                "slide_number": slide_num,
                "success": True,
                "applied_position": reference_position
            })

        except Exception as e:
            results.append({
                "slide_number": slide_num,
                "success": False,
                "error": str(e)
            })

    return {"data": "Success"}


def align_footnotes_to_reference(slide_map: Dict[str, Any], reference_slide_num: int, target_slide_numbers: list):
    """
    Align footnote positions of multiple slides to match a reference slide.
    
    Footnotes are detected as text boxes in the bottom 15% of the slide (top > slide_height * 0.85)
    
    This is a utility function, not an API endpoint.
    """
    # Get slide_ids_by_index
    slide_ids_by_index = list(slide_map.keys())
    total_slides = len(slide_ids_by_index)
    
    # Validate reference slide number
    ref_index = reference_slide_num - 1
    if ref_index < 0 or ref_index >= total_slides:
        return {
            "error": f"Invalid reference slide number. Valid range is 1-{total_slides}"
        }
    
    # Get reference slide
    ref_slide_id = slide_ids_by_index[ref_index]
    ref_slide = slide_map[ref_slide_id]
    
    # Standard slide dimensions in EMUs (PowerPoint uses 914400 EMUs per inch)
    # Standard slide is 10 inches wide x 7.5 inches tall
    slide_height = Inches(7.5)
    footnote_threshold = slide_height * 0.85
    
    # Find footnote in reference slide (text box in bottom 15% of slide)
    ref_footnote_shape = None
    
    for shape in ref_slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX or (
            shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE and shape.has_text_frame
        ):
            # Check if it's in the bottom 15% of the slide
            if shape.top > footnote_threshold:
                ref_footnote_shape = shape
                break  # Take the first footnote found
    
    if ref_footnote_shape is None:
        return {
            "error": f"Reference slide {reference_slide_num} does not have a footnote (text box in bottom 15% of slide)"
        }
    
    # Get reference position
    reference_position = {
        "left": ref_footnote_shape.left.inches,
        "top": ref_footnote_shape.top.inches,
        "width": ref_footnote_shape.width.inches,
        "height": ref_footnote_shape.height.inches
    }
    
    # Apply to target slides
    results = []
    
    for slide_num in target_slide_numbers:
        slide_index = slide_num - 1
        
        if slide_index < 0 or slide_index >= total_slides:
            results.append({
                "slide_number": slide_num,
                "success": False,
                "error": f"Invalid slide number. Valid range is 1-{total_slides}"
            })
            continue
        
        slide_id = slide_ids_by_index[slide_index]
        slide = slide_map[slide_id]
        
        # Find footnote shape (text box in bottom 15%)
        footnote_shape = None
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX or (
                shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE and shape.has_text_frame
            ):
                if shape.top > footnote_threshold:
                    footnote_shape = shape
                    break
        
        if footnote_shape is None:
            results.append({
                "slide_number": slide_num,
                "success": False,
                "error": "No footnote shape found in this slide (no text box in bottom 15%)"
            })
            continue
        
        try:
            footnote_shape.left = Inches(reference_position["left"])
            footnote_shape.top = Inches(reference_position["top"])
            footnote_shape.width = Inches(reference_position["width"])
            footnote_shape.height = Inches(reference_position["height"])
            
            results.append({
                "slide_number": slide_num,
                "success": True,
                "applied_position": reference_position
            })
        except Exception as e:
            results.append({
                "slide_number": slide_num,
                "success": False,
                "error": str(e)
            })
    
    return {"data": "Success"}


def get_title_coordinates(slide_map: Dict[str, Any], slide_id: str):
    """
    Get the position and dimensions of the title shape in a slide.
    
    This is a utility function, not an API endpoint.
    """
    if slide_id not in slide_map:
        return {"error": "slide not found"}
    
    slide = slide_map[slide_id]
    
    # Try to find the title shape
    title_shape = None
    
    # First check if slide has a title attribute
    if hasattr(slide.shapes, "title"):
        try:
            title_shape = slide.shapes.title
        except:
            pass
    
    # If not found, look for placeholder with type 1 (title)
    if title_shape is None:
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
                if hasattr(shape, "placeholder_format"):
                    if shape.placeholder_format.type == 1:
                        title_shape = shape
                        break
    
    if title_shape is None:
        return {
            "slide_id": slide_id,
            "has_title": False,
            "error": "No title shape found in this slide"
        }
    
    return {
        "slide_id": slide_id,
        "has_title": True,
        "title_text": title_shape.text if title_shape.has_text_frame else "",
        "coordinates": {
            "left": title_shape.left.inches,
            "top": title_shape.top.inches,
            "width": title_shape.width.inches,
            "height": title_shape.height.inches
        }
    }


def set_bulk_title_positions(slide_map: Dict[str, Any], request: dict):
    """
    Set the same title position for multiple slides.
    
    Expected request format:
    {
        "slide_numbers": [1, 2, 3, 5],  # 1-based slide numbers
        "position": {
            "left": 1.0,
            "top": 0.5,
            "width": 8.0,
            "height": 1.0
        }
    }
    
    This is a utility function, not an API endpoint.
    """
    slide_numbers = request["slide_numbers"]
    position = request["position"]
    
    # Convert slide_map (UUID keys) to index-based lookup
    slide_ids_by_index = list(slide_map.keys())  # Ordered list of UUIDs
    
    results = []
    
    for slide_num in slide_numbers:
        slide_index = slide_num - 1  # Convert 1-based to 0-based
        
        # Check if index is valid
        if slide_index < 0 or slide_index >= len(slide_ids_by_index):
            results.append({
                "slide_number": slide_num,
                "slide_index": slide_index,
                "success": False,
                "error": f"Invalid slide number. Valid range is 1-{len(slide_ids_by_index)}"
            })
            continue
        
        # Get the actual UUID for this slide index
        slide_id = slide_ids_by_index[slide_index]
        slide = slide_map[slide_id]
        
        # Find the title shape
        title_shape = None
        
        if hasattr(slide.shapes, "title"):
            try:
                title_shape = slide.shapes.title
            except:
                pass
        
        if title_shape is None:
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
                    if hasattr(shape, "placeholder_format"):
                        if shape.placeholder_format.type == 1:
                            title_shape = shape
                            break
        
        if title_shape is None:
            results.append({
                "slide_number": slide_num,
                "slide_id": slide_id,
                "success": False,
                "error": "No title shape found in this slide"
            })
            continue
        
        try:
            if "left" in position:
                title_shape.left = Inches(position["left"])
            if "top" in position:
                title_shape.top = Inches(position["top"])
            if "width" in position:
                title_shape.width = Inches(position["width"])
            if "height" in position:
                title_shape.height = Inches(position["height"])
            
            results.append({
                "slide_number": slide_num,
                "slide_id": slide_id,
                "success": True,
                "new_position": {
                    "left": title_shape.left.inches,
                    "top": title_shape.top.inches,
                    "width": title_shape.width.inches,
                    "height": title_shape.height.inches
                }
            })
        except Exception as e:
            results.append({
                "slide_number": slide_num,
                "slide_id": slide_id,
                "success": False,
                "error": str(e)
            })
    
    return {
        "total_slides": len(slide_numbers),
        "successful": sum(1 for r in results if r.get("success")),
        "failed": sum(1 for r in results if not r.get("success")),
        "results": results
    }


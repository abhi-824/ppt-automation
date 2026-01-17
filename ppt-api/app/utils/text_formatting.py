"""
Utility functions for text formatting operations.
"""
import re
from pptx.util import Pt, Inches
from pptx.enum.text import MSO_AUTO_SIZE, MSO_VERTICAL_ANCHOR


def apply_markdown_to_text_frame(text_frame, text: str):
    """
    Parses simple markdown syntax and applies formatting:
    - # Heading 1
    - ## Heading 2
    - **bold**
    - *italic*
    """
    text_frame.clear()

    for line in text.split('\n'):
        line = line.strip()
        if not line:
            continue

        p = text_frame.add_paragraph()

        # Headings
        if line.startswith('## '):
            run = p.add_run()
            run.text = line[3:]
            run.font.size = Pt(28)
            run.font.bold = True
        elif line.startswith('# '):
            run = p.add_run()
            run.text = line[2:]
            run.font.size = Pt(36)
            run.font.bold = True
        else:
            # Bold/Italic inline formatting
            cursor = 0
            for match in re.finditer(r'(\*\*.*?\*\*|\*.*?\*)', line):
                start, end = match.span()
                if start > cursor:
                    run = p.add_run()
                    run.text = line[cursor:start]

                md = match.group()
                run = p.add_run()
                if md.startswith('**'):
                    run.text = md[2:-2]
                    run.font.bold = True
                elif md.startswith('*'):
                    run.text = md[1:-1]
                    run.font.italic = True
                cursor = end

            if cursor < len(line):
                run = p.add_run()
                run.text = line[cursor:]


def configure_textbox_frame(text_frame):
    """Configure text box frame settings."""
    text_frame.word_wrap = True
    text_frame.auto_size = MSO_AUTO_SIZE.NONE
    text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
    text_frame.margin_top = Inches(0.05)
    text_frame.margin_bottom = Inches(0.05)
    text_frame.margin_left = Inches(0.05)
    text_frame.margin_right = Inches(0.05)


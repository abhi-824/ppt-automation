from fastapi import FastAPI, UploadFile, File
from pydantic import BaseModel
from typing import List, Dict, Optional, Any
from pptx import Presentation
from pptx.util import Inches
from pptx.enum.text import MSO_AUTO_SIZE, MSO_VERTICAL_ANCHOR
import uuid
import os
import re
import io
import base64
from pptx.util import Pt
from fastapi.middleware.cors import CORSMiddleware
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.shapes import PP_PLACEHOLDER

from .themes.theme import Theme, THEMES
from .components.layouts import (
    HeaderWithImage, BulletWithTitle, TwoColumnText,
    ComparisonTable, IconList, QuoteBlock, Timeline, ProcessFlow, StatisticHighlight, CalloutBox, SectionDivider
)

app = FastAPI()
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)
prs = Presentation()
slide_map: Dict[str, object] = {}
current_theme: Theme = THEMES["default"]

# Models
class ComponentContent(BaseModel):
    component_type: str
    content: Dict[str, Any]

class ThemeRequest(BaseModel):
    theme_name: str

class TextRequest(BaseModel):
    text: str
    left: Optional[float] = None
    top: Optional[float] = None
    width: Optional[float] = None
    height: Optional[float] = None

class BulletPointsRequest(BaseModel):
    points: List[str]
    left: Optional[float] = None
    top: Optional[float] = None
    width: Optional[float] = None
    height: Optional[float] = None

class TextBoxRequest(BaseModel):
    text: str
    left: float = 1.0
    top: float = 1.0
    width: float = 4.0
    height: float = 1.0

class SlideContentRequest(BaseModel):
    title: str
    content: str

class HeaderWithImageContent(BaseModel):
    title: str
    image_path: str
    left: Optional[float] = 1.0
    top: Optional[float] = 2.0
    width: Optional[float] = 8.0
    height: Optional[float] = 4.0

class TwoColumnTextContent(BaseModel):
    title: str
    left_text: str
    right_text: str
    left: Optional[float] = 1.0
    top: Optional[float] = 1.0
    width: Optional[float] = 4.0
    height: Optional[float] = 4.0
    column_gap: Optional[float] = 0.5

class SlideBase64Request(BaseModel):
    slideBase64: str

@app.post("/set/slideBase64")
def set_slide_base64(req: SlideBase64Request):
    """
    Initialize the presentation by clearing prs and loading from base64.
    This sets the current PowerPoint presentation we are working with.
    """
    global prs, slide_map
    
    try:
        # Clear the current presentation and slide map
        prs = Presentation()
        slide_map.clear()
        
        # Decode the base64 string
        decoded_bytes = base64.b64decode(req.slideBase64)
        
        print("hi")
        # Load the presentation from the decoded bytes
        buffer = io.BytesIO(decoded_bytes)
        print("hi2")
        prs = Presentation(buffer)
        print("hi3")
        # Optionally, populate slide_map with existing slides
        for i, slide in enumerate(prs.slides):
            slide_id = str(uuid.uuid4())
            slide_map[slide_id] = slide
        print("slide_map", slide_map)
        print(len(prs.slides))
        print(list[str](slide_map.keys()))
        return {
            "status": "ok",
            "message": "Presentation initialized from base64",
            "slide_count": len(prs.slides),
            "slide_ids": list(slide_map.keys())
        }
    except Exception as e:
        print("error", e)
        return {
            "status": "error",
            "error": str(e)
        }
# Theme endpoints
@app.get("/themes")
def list_themes():
    return {"themes": list(THEMES.keys())}

@app.post("/theme")
def set_theme(req: ThemeRequest):
    global current_theme
    if req.theme_name not in THEMES:
        return {"error": "theme not found"}
    current_theme = THEMES[req.theme_name]
    return {"status": "ok"}

# Individual slide element endpoints - used by MCP server
@app.post("/slide/{slide_id}/title")
def add_title(slide_id: str, req: TextRequest):
    """Add title to a slide - used by MCP server"""
    if slide_id not in slide_map:
        return {"error": "slide not found"}
    
    slide = slide_map[slide_id]
    if not slide.shapes.title:
        # If no title placeholder, create a text box
        if req.left is not None and req.top is not None:
            text_box = slide.shapes.add_textbox(
                Inches(req.left), Inches(req.top),
                Inches(req.width or 8), Inches(req.height or 1.5)
            )
            configure_textbox_frame(text_box.text_frame)
            apply_markdown_to_text_frame(text_box.text_frame, req.text)
            if current_theme:
                current_theme.get_style("title").apply_to_text_frame(text_box.text_frame)
        else:
            return {"error": "position required when no title placeholder exists"}
    else:
        configure_textbox_frame(slide.shapes.title.text_frame)
        apply_markdown_to_text_frame(slide.shapes.title.text_frame, req.text)
        if current_theme:
            current_theme.get_style("title").apply_to_text_frame(slide.shapes.title.text_frame)
    return {"status": "ok"}

@app.post("/slide/{slide_id}/subtitle")
def add_subtitle(slide_id: str, req: TextRequest):
    """Add subtitle to a slide - used by MCP server"""
    if slide_id not in slide_map:
        return {"error": "slide not found"}
    
    slide = slide_map[slide_id]
    # Find subtitle placeholder (usually index 1)
    subtitle = None
    for shape in slide.shapes:
        if hasattr(shape, "placeholder_format"):
            if shape.placeholder_format.idx == 1:
                subtitle = shape
                break
    
    if not subtitle:
        # If no subtitle placeholder, create a text box
        if req.left is not None and req.top is not None:
            text_box = slide.shapes.add_textbox(
                Inches(req.left), Inches(req.top),
                Inches(req.width or 8), Inches(req.height or 1)
            )
            configure_textbox_frame(text_box.text_frame)
            apply_markdown_to_text_frame(text_box.text_frame, req.text)
            if current_theme:
                current_theme.get_style("subtitle").apply_to_text_frame(text_box.text_frame)
        else:
            return {"error": "position required when no subtitle placeholder exists"}
    else:
        configure_textbox_frame(subtitle.text_frame)
        apply_markdown_to_text_frame(subtitle.text_frame, req.text)
        if current_theme:
            current_theme.get_style("subtitle").apply_to_text_frame(subtitle.text_frame)
    return {"status": "ok"}

@app.post("/slide/{slide_id}/bullet_points")
def add_bullet_points(slide_id: str, req: BulletPointsRequest):
    """Add bullet points to a slide - used by MCP server"""
    if slide_id not in slide_map:
        return {"error": "slide not found"}
    
    slide = slide_map[slide_id]
    
    left = req.left if req.left is not None else 1
    top = req.top if req.top is not None else 2
    width = req.width if req.width is not None else 8
    height = req.height if req.height is not None else 4

    shape = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = shape.text_frame
    configure_textbox_frame(tf)
    
    markdown_text = "\n".join(req.points)  # Join bullet points as lines
    apply_markdown_to_text_frame(tf, markdown_text)  # ✅ Now it's correct

    if current_theme:
        current_theme.get_style("bullet").apply_to_text_frame(tf)

    return {"status": "ok"}

@app.post("/slide/{slide_id}/text_box")
def add_text_box(slide_id: str, req: TextBoxRequest):
    """Add a text box to a slide - used by MCP server"""
    if slide_id not in slide_map:
        return {"error": "slide not found"}
    
    slide = slide_map[slide_id]
    left = req.left if req.left is not None else 1
    top = req.top if req.top is not None else 1
    width = req.width if req.width is not None else 8
    height = req.height if req.height is not None else 1
    text_box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    configure_textbox_frame(text_box.text_frame)
    apply_markdown_to_text_frame(text_box.text_frame, req.text)
    if current_theme:
        current_theme.get_style("body").apply_to_text_frame(text_box.text_frame)
    return {"status": "ok"}

# mponentmponent endpoints
@app.post("/slide/{slide_id}/component")
def add_component(slide_id: str, req: ComponentContent):
    if slide_id not in slide_map:
        return {"error": "slide not found"}
    
    slide = slide_map[slide_id]
    component_map = {
        "header_with_image": HeaderWithImage,
        "bullet_with_title": BulletWithTitle,
        "two_column_text": TwoColumnText,
        "comparison_table": ComparisonTable,
        "icon_list": IconList,
        "quote_block": QuoteBlock,
        "timeline": Timeline,
        "process_flow": ProcessFlow,
        "statistic_highlight": StatisticHighlight,
        "callout_box": CalloutBox,
        "section_divider": SectionDivider
    }
    
    if req.component_type not in component_map:
        return {"error": "invalid component type"}
    
    component_class = component_map[req.component_type]
    component = component_class(theme=current_theme)
    component.render(slide, req.content)
    
    return {"status": "ok"}

# Basic slide management
@app.post("/slide")
def create_slide(layout: int = 0):
    """Create a new blank slide - used by MCP server"""
    slide = prs.slides.add_slide(prs.slide_layouts[layout])
    slide_id = str(uuid.uuid4())
    slide_map[slide_id] = slide
    
    # Apply current theme to new slide
    current_theme.apply_to_slide(slide)
    
    return {"slide_id": slide_id}

@app.post("/slide/blank")
def create_blank_slide():
    """Create a new blank slide without any placeholders"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Layout 6 is blank
    slide_id = str(uuid.uuid4())
    slide_map[slide_id] = slide
    
    # Apply current theme to new slide
    current_theme.apply_to_slide(slide)
    
    return {"slide_id": slide_id}

@app.delete("/slide/{slide_id}")
def delete_slide(slide_id: str):
    """Delete a slide - used by MCP server"""
    if slide_id in slide_map:
        slide = slide_map.pop(slide_id)
        # Find the slide in the presentation and remove it
        for idx, sld in enumerate(prs.slides):
            if sld == slide:
                xml_slides = prs.slides._sldIdLst
                xml_slides.remove(xml_slides[idx])
                return {"status": "deleted"}
        return {"error": "slide not found in presentation"}
    return {"error": "slide not found"}

@app.get("/slides")
def list_slides():
    """List all slides - used by MCP server"""
    return {"slide_ids": list(slide_map.keys())}

@app.post("/save")
def save_presentation(filename: str = "output.pptx"):
    """Save the presentation - used by MCP server"""
    prs.save(filename)
    return {"status": "saved", "filename": filename}

# Add this endpoint to your FastAPI app
@app.get("/presentation/base64")
def get_presentation_base64():
    """
    Export the current presentation as base64 string
    Returns: {"base64": "...", "filename": "presentation.pptx"}
    """
    try:
        # Save presentation to BytesIO buffer
        buffer = io.BytesIO()
        prs.save(buffer)
        buffer.seek(0)
        
        # Convert to base64
        base64_data = base64.b64encode(buffer.read()).decode('utf-8')
        
        return {
            "status": "ok",
            "base64": base64_data,
            "filename": "presentation.pptx",
            "slide_count": len(prs.slides)
        }
    except Exception as e:
        return JSONResponse(
            status_code=500,
            content={"error": str(e)}
        )

@app.get("/presentation/preview")
def get_presentation_preview():
    """
    Get presentation metadata and base64 for preview
    """
    try:
        buffer = io.BytesIO()
        prs.save(buffer)
        buffer.seek(0)
        
        base64_data = base64.b64encode(buffer.read()).decode('utf-8')
        
        return {
            "status": "ok",
            "base64": base64_data,
            "slide_count": len(prs.slides),
            "slide_ids": list(slide_map.keys()),
            "current_theme": current_theme.name if hasattr(current_theme, 'name') else "default"
        }
    except Exception as e:
        return JSONResponse(
            status_code=500,
            content={"error": str(e)}
        )

@app.post("/presentation/reset")
def reset_presentation():
    """Reset the presentation to start fresh"""
    global prs, slide_map
    prs = Presentation()
    slide_map = {}
    return {"status": "ok", "message": "Presentation reset"}


def apply_markdown_to_text_frame(text_frame, text: str):
    """
    Parses simple markdown syntax and applies formatting:
    - # Heading 1
    - ## Heading 2
    - **bold**
    - *italic*
    """
    import re
    text_frame.clear()

    for line in text.split('\n'):
        line = line.strip()
        if not line:
            continue

        p = text_frame.add_paragraph()

        # Headings
        if line.startswith('## '):
            run = p.add_run()
            run.text = line[3:]
            run.font.size = Pt(28)
            run.font.bold = True
        elif line.startswith('# '):
            run = p.add_run()
            run.text = line[2:]
            run.font.size = Pt(36)
            run.font.bold = True
        else:
            # Bold/Italic inline formatting
            cursor = 0
            for match in re.finditer(r'(\*\*.*?\*\*|\*.*?\*)', line):
                start, end = match.span()
                if start > cursor:
                    run = p.add_run()
                    run.text = line[cursor:start]

                md = match.group()
                run = p.add_run()
                if md.startswith('**'):
                    run.text = md[2:-2]
                    run.font.bold = True
                elif md.startswith('*'):
                    run.text = md[1:-1]
                    run.font.italic = True
                cursor = end

            if cursor < len(line):
                run = p.add_run()
                run.text = line[cursor:]

def configure_textbox_frame(text_frame):
    text_frame.word_wrap = True
    text_frame.auto_size = MSO_AUTO_SIZE.NONE
    text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
    text_frame.margin_top = Inches(0.05)
    text_frame.margin_bottom = Inches(0.05)
    text_frame.margin_left = Inches(0.05)
    text_frame.margin_right = Inches(0.05)

# API 1: Get all shapes from a slide
@app.get("/slide/{slide_id}/shapes")
def get_slide_shapes(slide_id: str):
    """
    Get detailed information about all shapes in a slide including type, position, and content
    """
    if slide_id not in slide_map:
        return {"error": "slide not found"}
    
    slide = slide_map[slide_id]
    shapes_info = []
    
    for idx, shape in enumerate(slide.shapes):
        shape_info = {
            "shape_id": f"shape_{idx}",
            "shape_index": idx,
            "shape_name": shape.name,
            "has_text": shape.has_text_frame,
            "text_content": "",
            "position": {
                "left": shape.left.inches,
                "top": shape.top.inches,
                "width": shape.width.inches,
                "height": shape.height.inches
            },
            "is_title": False,
            "is_subtitle": False,
            "shape_type": "unknown"
        }
        
        # Determine shape type
        if shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
            shape_info["shape_type"] = "placeholder"
            # Check if it's title or subtitle
            if hasattr(shape, "placeholder_format"):
                if shape.placeholder_format.type == 1:  # Title
                    shape_info["is_title"] = True
                    shape_info["shape_type"] = "title"
                elif shape.placeholder_format.type == 2:  # Subtitle/Body
                    shape_info["is_subtitle"] = True
                    shape_info["shape_type"] = "subtitle"
        elif shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
            shape_info["shape_type"] = "text_box"
        elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            shape_info["shape_type"] = "picture"
        elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            shape_info["shape_type"] = "table"
        elif shape.shape_type == MSO_SHAPE_TYPE.CHART:
            shape_info["shape_type"] = "chart"
        elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            shape_info["shape_type"] = "group"
        elif shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            shape_info["shape_type"] = "auto_shape"
        
        # Also check if it's the title by checking if it's the title shape
        if hasattr(slide.shapes, "title") and slide.shapes.title == shape:
            shape_info["is_title"] = True
            shape_info["shape_type"] = "title"
        
        # Get text content if available
        if shape.has_text_frame:
            text_parts = []
            for paragraph in shape.text_frame.paragraphs:
                text_parts.append(paragraph.text)
            shape_info["text_content"] = "\n".join(text_parts)
        
        shapes_info.append(shape_info)
    
    return {
        "slide_id": slide_id,
        "shape_count": len(shapes_info),
        "shapes": shapes_info
    }


# Model for setting title position
class TitlePositionRequest(BaseModel):
    left: float
    top: float
    width: float
    height: float


# API 2: Set title position on a slide
@app.post("/slide/{slide_id}/title/position")
def set_title_position(slide_id: str, req: TitlePositionRequest):
    """
    Set the position (left, top, width, height) of the title shape on a slide
    """
    if slide_id not in slide_map:
        return {"error": "slide not found"}
    
    slide = slide_map[slide_id]
    
    # Find the title shape
    title_shape = None
    
    # Method 1: Try the standard title property
    if hasattr(slide.shapes, "title") and slide.shapes.title:
        title_shape = slide.shapes.title
    else:
        # Method 2: Look for placeholder with title type
        for shape in slide.shapes:
            if hasattr(shape, "placeholder_format"):
                if shape.placeholder_format.type == 1:  # Title placeholder
                    title_shape = shape
                    break
    
    if not title_shape:
        return {
            "error": "title shape not found on this slide",
            "message": "This slide may not have a title placeholder"
        }
    
    # Update position
    title_shape.left = Inches(req.left)
    title_shape.top = Inches(req.top)
    title_shape.width = Inches(req.width)
    title_shape.height = Inches(req.height)
    
    return {
        "status": "ok",
        "message": "Title position updated successfully",
        "slide_id": slide_id,
        "new_position": {
            "left": req.left,
            "top": req.top,
            "width": req.width,
            "height": req.height
        }
    }
@app.get("/slide/{slide_id}/title/coordinates")
def get_title_coordinates(slide_id: str):
    """
    Get the position and dimensions of the title shape in a slide
    """
    if slide_id not in slide_map:
        return {"error": "slide not found"}
    
    slide = slide_map[slide_id]
    
    # Try to find the title shape
    title_shape = None
    
    # First check if slide has a title attribute
    if hasattr(slide.shapes, "title"):
        try:
            title_shape = slide.shapes.title
        except:
            pass
    
    # If not found, look for placeholder with type 1 (title)
    if title_shape is None:
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
                if hasattr(shape, "placeholder_format"):
                    if shape.placeholder_format.type == 1:
                        title_shape = shape
                        break
    
    if title_shape is None:
        return {
            "slide_id": slide_id,
            "has_title": False,
            "error": "No title shape found in this slide"
        }
    
    return {
        "slide_id": slide_id,
        "has_title": True,
        "title_text": title_shape.text if title_shape.has_text_frame else "",
        "coordinates": {
            "left": title_shape.left.inches,
            "top": title_shape.top.inches,
            "width": title_shape.width.inches,
            "height": title_shape.height.inches
        }
    }

@app.post("/slides/bulk/title/position")
def set_bulk_title_positions(request: dict):
    """
    Set the same title position for multiple slides
    
    Expected body format:
    {
        "slide_numbers": [1, 2, 3, 5],  # 1-based slide numbers
        "position": {
            "left": 1.0,
            "top": 0.5,
            "width": 8.0,
            "height": 1.0
        }
    }
    
    Position values are in inches and will be applied to all specified slides
    """
    print(slide_map)

    if "slide_numbers" not in request or "position" not in request:
        return {"error": "Missing 'slide_numbers' or 'position' field in request body"}
    
    slide_numbers = request["slide_numbers"]
    position = request["position"]
    
    if not isinstance(slide_numbers, list) or len(slide_numbers) == 0:
        return {"error": "slide_numbers must be a non-empty array"}
    
    # Convert slide_map (UUID keys) to index-based lookup
    # slide_map keys are UUIDs, but we need to access slides by their index
    slide_ids_by_index = list(slide_map.keys())  # Ordered list of UUIDs
    
    results = []
    
    for slide_num in slide_numbers:
        slide_index = slide_num - 1  # Convert 1-based to 0-based
        
        # Check if index is valid
        if slide_index < 0 or slide_index >= len(slide_ids_by_index):
            results.append({
                "slide_number": slide_num,
                "slide_index": slide_index,
                "success": False,
                "error": f"Invalid slide number. Valid range is 1-{len(slide_ids_by_index)}"
            })
            continue
        
        # Get the actual UUID for this slide index
        slide_id = slide_ids_by_index[slide_index]
        slide = slide_map[slide_id]
        
        # Find the title shape
        title_shape = None
        
        if hasattr(slide.shapes, "title"):
            try:
                title_shape = slide.shapes.title
            except:
                pass
        
        if title_shape is None:
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
                    if hasattr(shape, "placeholder_format"):
                        if shape.placeholder_format.type == 1:
                            title_shape = shape
                            break
        
        if title_shape is None:
            results.append({
                "slide_number": slide_num,
                "slide_id": slide_id,
                "success": False,
                "error": "No title shape found in this slide"
            })
            continue
        
        try:
            # Set the position
            from pptx.util import Inches
            
            if "left" in position:
                title_shape.left = Inches(position["left"])
            if "top" in position:
                title_shape.top = Inches(position["top"])
            if "width" in position:
                title_shape.width = Inches(position["width"])
            if "height" in position:
                title_shape.height = Inches(position["height"])
            
            results.append({
                "slide_number": slide_num,
                "slide_id": slide_id,
                "success": True,
                "new_position": {
                    "left": title_shape.left.inches,
                    "top": title_shape.top.inches,
                    "width": title_shape.width.inches,
                    "height": title_shape.height.inches
                }
            })
        except Exception as e:
            results.append({
                "slide_number": slide_num,
                "slide_id": slide_id,
                "success": False,
                "error": str(e)
            })
    
    return {
        "total_slides": len(slide_numbers),
        "successful": sum(1 for r in results if r.get("success")),
        "failed": sum(1 for r in results if not r.get("success")),
        "results": results
    }

@app.post("/slides/align_titles_to_reference")
def align_titles_to_reference(request: dict):
    """
    Align title positions of multiple slides to match a reference slide
    
    Expected body format:
    {
        "reference_slide_number": 2,  # 1-based slide number to copy from
        "target_slide_numbers": [1, 3, 4, 5]  # 1-based slide numbers to update
    }
    
    This internally uses:
    1. GET /slide/{slide_id}/title/coordinates - to get reference position
    2. POST /slides/bulk/title/position - to apply to targets
    """
    
    if "reference_slide_number" not in request or "target_slide_numbers" not in request:
        return {"error": "Missing 'reference_slide_number' or 'target_slide_numbers' field"}
    
    reference_slide_num = request["reference_slide_number"]
    target_slide_numbers = request["target_slide_numbers"]
    
    if not isinstance(target_slide_numbers, list) or len(target_slide_numbers) == 0:
        return {"error": "target_slide_numbers must be a non-empty array"}
    
    # Get slide_ids_by_index
    slide_ids_by_index = list(slide_map.keys())
    total_slides = len(slide_ids_by_index)
    
    # Validate reference slide number
    ref_index = reference_slide_num - 1
    if ref_index < 0 or ref_index >= total_slides:
        return {
            "error": f"Invalid reference slide number. Valid range is 1-{total_slides}"
        }
    
    # Step 1: Get reference slide coordinates using existing API logic
    ref_slide_id = slide_ids_by_index[ref_index]
    ref_coordinates_response = get_title_coordinates(ref_slide_id)
    
    # Check if reference slide has a title
    if "error" in ref_coordinates_response or not ref_coordinates_response.get("has_title"):
        print("error", reference_slide_num)
        return {
            "error": f"Reference slide {reference_slide_num} does not have a title shape",
            "details": ref_coordinates_response
        }
    
    # Extract position from reference
    reference_position = ref_coordinates_response["coordinates"]
    
    # Step 2: Apply position to target slides using existing bulk API logic
    bulk_request = {
        "slide_numbers": target_slide_numbers,
        "position": reference_position
    }
    
    bulk_response = set_bulk_title_positions(bulk_request)
    
    # Combine and return
    return {"data":"Success"}

def _find_subtitle(shapes):
    """
    Find a subtitle shape using BOTH:
    - Official python-pptx placeholder type (PP_PLACEHOLDER.SUBTITLE)
    - User-defined working logic: placeholder_format.type == 2
    """
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
            if hasattr(shape, "placeholder_format"):
                ph_type = shape.placeholder_format.type

                # Official subtitle placeholder
                if ph_type == PP_PLACEHOLDER.SUBTITLE:
                    return shape

                # User custom logic (your working system)
                # type == 2 means Subtitle/Body for your template
                if ph_type == 2:
                    return shape
    return None


@app.post("/slides/align_subtitles_to_reference")
def align_subtitles_to_reference(request: dict):
    """
    Align subtitle positions of multiple slides to match a reference slide.
    """

    if "reference_slide_number" not in request or "target_slide_numbers" not in request:
        return {"error": "Missing 'reference_slide_number' or 'target_slide_numbers' field"}

    reference_slide_num = request["reference_slide_number"]
    target_slide_numbers = request["target_slide_numbers"]

    if not isinstance(target_slide_numbers, list) or len(target_slide_numbers) == 0:
        return {"error": "target_slide_numbers must be a non-empty array"}

    # Slide index handling
    slide_ids_by_index = list(slide_map.keys())
    total_slides = len(slide_ids_by_index)

    ref_index = reference_slide_num - 1
    if ref_index < 0 or ref_index >= total_slides:
        return {
            "error": f"Invalid reference slide number. Valid range is 1-{total_slides}"
        }

    # Get the reference slide object
    ref_slide = slide_map[slide_ids_by_index[ref_index]]

    # Find reference subtitle
    ref_subtitle_shape = _find_subtitle(ref_slide.shapes)

    if ref_subtitle_shape is None:
        return {
            "error": f"Reference slide {reference_slide_num} does not contain a subtitle placeholder"
        }

    # Extract reference subtitle geometry
    reference_position = {
        "left": ref_subtitle_shape.left.inches,
        "top": ref_subtitle_shape.top.inches,
        "width": ref_subtitle_shape.width.inches,
        "height": ref_subtitle_shape.height.inches
    }

    # Apply to targets
    results = []

    for slide_num in target_slide_numbers:
        slide_index = slide_num - 1

        if slide_index < 0 or slide_index >= total_slides:
            results.append({
                "slide_number": slide_num,
                "success": False,
                "error": f"Invalid slide number. Valid range is 1-{total_slides}"
            })
            continue

        slide = slide_map[slide_ids_by_index[slide_index]]

        subtitle_shape = _find_subtitle(slide.shapes)

        if subtitle_shape is None:
            results.append({
                "slide_number": slide_num,
                "success": False,
                "error": "No subtitle placeholder found on this slide"
            })

            continue

        try:
            subtitle_shape.left = Inches(reference_position["left"])
            subtitle_shape.top = Inches(reference_position["top"])
            subtitle_shape.width = Inches(reference_position["width"])
            subtitle_shape.height = Inches(reference_position["height"])

            results.append({
                "slide_number": slide_num,
                "success": True,
                "applied_position": reference_position
            })

        except Exception as e:
            results.append({
                "slide_number": slide_num,
                "success": False,
                "error": str(e)
            })

    return {"data":"Success"}

@app.post("/slides/align_footnotes_to_reference")
def align_footnotes_to_reference(request: dict):
    """
    Align footnote positions of multiple slides to match a reference slide
    
    Footnotes are detected as text boxes in the bottom 15% of the slide (top > slide_height * 0.85)
    
    Expected body format:
    {
        "reference_slide_number": 2,  # 1-based slide number to copy from
        "target_slide_numbers": [1, 3, 4, 5]  # 1-based slide numbers to update
    }
    """
    
    if "reference_slide_number" not in request or "target_slide_numbers" not in request:
        return {"error": "Missing 'reference_slide_number' or 'target_slide_numbers' field"}
    
    reference_slide_num = request["reference_slide_number"]
    target_slide_numbers = request["target_slide_numbers"]
    
    if not isinstance(target_slide_numbers, list) or len(target_slide_numbers) == 0:
        return {"error": "target_slide_numbers must be a non-empty array"}
    
    # Get slide_ids_by_index
    slide_ids_by_index = list(slide_map.keys())
    total_slides = len(slide_ids_by_index)
    
    # Validate reference slide number
    ref_index = reference_slide_num - 1
    if ref_index < 0 or ref_index >= total_slides:
        return {
            "error": f"Invalid reference slide number. Valid range is 1-{total_slides}"
        }
    
    # Get reference slide
    ref_slide_id = slide_ids_by_index[ref_index]
    ref_slide = slide_map[ref_slide_id]
    
    # Standard slide dimensions in EMUs (PowerPoint uses 914400 EMUs per inch)
    # Standard slide is 10 inches wide x 7.5 inches tall
    slide_height = Inches(7.5)
    footnote_threshold = slide_height * 0.85
    
    # Find footnote in reference slide (text box in bottom 15% of slide)
    ref_footnote_shape = None
    
    for shape in ref_slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX or (
            shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE and shape.has_text_frame
        ):
            # Check if it's in the bottom 15% of the slide
            if shape.top > footnote_threshold:
                ref_footnote_shape = shape
                break  # Take the first footnote found
    
    if ref_footnote_shape is None:
        return {
            "error": f"Reference slide {reference_slide_num} does not have a footnote (text box in bottom 15% of slide)"
        }
    
    # Get reference position
    reference_position = {
        "left": ref_footnote_shape.left.inches,
        "top": ref_footnote_shape.top.inches,
        "width": ref_footnote_shape.width.inches,
        "height": ref_footnote_shape.height.inches
    }
    
    # Apply to target slides
    results = []
    
    for slide_num in target_slide_numbers:
        slide_index = slide_num - 1
        
        if slide_index < 0 or slide_index >= total_slides:
            results.append({
                "slide_number": slide_num,
                "success": False,
                "error": f"Invalid slide number. Valid range is 1-{total_slides}"
            })
            continue
        
        slide_id = slide_ids_by_index[slide_index]
        slide = slide_map[slide_id]
        
        # Find footnote shape (text box in bottom 15%)
        footnote_shape = None
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX or (
                shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE and shape.has_text_frame
            ):
                if shape.top > footnote_threshold:
                    footnote_shape = shape
                    break
        
        if footnote_shape is None:
            results.append({
                "slide_number": slide_num,
                "success": False,
                "error": "No footnote shape found in this slide (no text box in bottom 15%)"
            })
            continue
        
        try:
            footnote_shape.left = Inches(reference_position["left"])
            footnote_shape.top = Inches(reference_position["top"])
            footnote_shape.width = Inches(reference_position["width"])
            footnote_shape.height = Inches(reference_position["height"])
            
            results.append({
                "slide_number": slide_num,
                "success": True,
                "applied_position": reference_position
            })
        except Exception as e:
            results.append({
                "slide_number": slide_num,
                "success": False,
                "error": str(e)
            })
    
    return {"data":"Success"}

@app.post("/slides/align_shapes_to_reference")
def align_shapes_to_reference(request: dict):
    """
    Align multiple shape types (title, subtitle, footnote) to match a reference slide
    
    Expected body format:
    {
        "reference_slide_number": 2,
        "target_slide_numbers": [1, 3, 4, 5],
        "shapes_to_align": ["title", "subtitle", "footnote"]
    }
    """
    
    if "reference_slide_number" not in request or "target_slide_numbers" not in request or "shapes_to_align" not in request:
        return {"error": "Missing required fields: 'reference_slide_number', 'target_slide_numbers', or 'shapes_to_align'"}
    
    reference_slide_num = request["reference_slide_number"]
    target_slide_numbers = request["target_slide_numbers"]
    shapes_to_align = request["shapes_to_align"]
    
    # Validate inputs
    if not isinstance(target_slide_numbers, list) or len(target_slide_numbers) == 0:
        return {"error": "target_slide_numbers must be a non-empty array"}
    
    if not isinstance(shapes_to_align, list) or len(shapes_to_align) == 0:
        return {"error": "shapes_to_align must be a non-empty array"}
    
    valid_shapes = ["title", "subtitle", "footnote"]
    for shape in shapes_to_align:
        if shape not in valid_shapes:
            return {"error": f"Invalid shape type '{shape}'. Valid types: {valid_shapes}"}
    
    # Dispatch to individual alignment functions
    results = {}
    
    if "title" in shapes_to_align:
        title_result = align_titles_to_reference({
            "reference_slide_number": reference_slide_num,
            "target_slide_numbers": target_slide_numbers
        })
        results["title"] = title_result
    
    if "subtitle" in shapes_to_align:
        subtitle_result = align_subtitles_to_reference({
            "reference_slide_number": reference_slide_num,
            "target_slide_numbers": target_slide_numbers
        })
        results["subtitle"] = subtitle_result
    
    if "footnote" in shapes_to_align:
        footnote_result = align_footnotes_to_reference({
            "reference_slide_number": reference_slide_num,
            "target_slide_numbers": target_slide_numbers
        })
        results["footnote"] = footnote_result
    
    return {
        "reference_slide": reference_slide_num,
        "target_slides": target_slide_numbers,
        "shapes_aligned": shapes_to_align,
        "results": results
    }
if __name__ == "__main__":
    main()

if __name__ == "__main__":
    import uvicorn
    uvicorn.run("app.main:app", host="127.0.0.1", port=8000, reload=True) 
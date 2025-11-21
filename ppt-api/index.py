import json
import subprocess
from typing import List, Dict, Optional
from pptx import Presentation
import re
from pptx.util import Inches, Pt
from fastapi import FastAPI, UploadFile, File
from pydantic import BaseModel
import uuid
import os

app = FastAPI()
prs = Presentation()
slide_map: Dict[str, object] = {}

# 1) Initialize
prs = Presentation()

class TitleRequest(BaseModel):
    text: str

class SubtitleRequest(BaseModel):
    text: str

class ParagraphRequest(BaseModel):
    text: str

class BulletPointsRequest(BaseModel):
    points: List[str]

class BackgroundRequest(BaseModel):
    color: str
    
class TextBoxRequest(BaseModel):
    text: str
    left: float = 1.0
    top: float = 1.0
    width: float = 4.0
    height: float = 1.0

class ImageRequest(BaseModel):
    image_path: str  # Or use UploadFile for direct upload
    left: Optional[float] = 1.0
    top: Optional[float] = 2.0
    width: Optional[float] = 5.0
    height: Optional[float] = 3.0

@app.delete("/slide/{slide_id}")
def delete_slide(slide_id: str):
    if slide_id in slide_map:
        slide = slide_map.pop(slide_id)
        prs.slides._sldIdLst.remove(slide._element.getparent())
        return {"status": "deleted"}
    return {"error": "slide not found"}

@app.get("/slides")
def list_slides():
    return {"slide_ids": list(slide_map.keys())}

@app.post("/slide/{slide_id}/bullet_points")
def add_bullets(slide_id: str, req: BulletPointsRequest):
    slide = slide_map[slide_id]
    tf = slide.shapes.placeholders[1].text_frame
    tf.clear()
    for pt in req.points:
        p = tf.add_paragraph()
        p.text = pt
    return {"status": "ok"}

@app.post("/slide")
def create_slide(layout: int = 0):
    slide = prs.slides.add_slide(prs.slide_layouts[layout])
    slide_id = str(uuid.uuid4())
    slide_map[slide_id] = slide
    return {"slide_id": slide_id}

@app.post("/slide/{slide_id}/title")
def add_title(slide_id: str, req: TitleRequest):
    slide = slide_map[slide_id]
    slide.shapes.title.text = req.text
    return {"status": "ok"}

@app.post("/slide/{slide_id}/subtitle")
def add_subtitle(slide_id: str, req: SubtitleRequest):
    slide = slide_map[slide_id]
    slide.placeholders[1].text = req.text
    return {"status": "ok"}

@app.post("/slide/{slide_id}/textbox")
def add_textbox(slide_id: str, req: TextBoxRequest):
    slide = slide_map[slide_id]
    txBox = slide.shapes.add_textbox(
        Inches(req.left), Inches(req.top), Inches(req.width), Inches(req.height)
    )
    txBox.text_frame.text = req.text
    return {"status": "ok"}

@app.post("/slide/{slide_id}/paragraph")
def add_paragraph(slide_id: str, req: ParagraphRequest):
    slide = slide_map[slide_id]
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(5), Inches(1))
    txBox.text_frame.text = req.text
    return {"status": "ok"}

@app.post("/slide/{slide_id}/image")
def add_image(slide_id: str, req: ImageRequest):
    slide = slide_map[slide_id]
    slide.shapes.add_picture(
        req.image_path,
        Inches(req.left),
        Inches(req.top),
        Inches(req.width),
        Inches(req.height)
    )
    return {"status": "ok"}

@app.post("/slide/{slide_id}/background")
def set_background(slide_id: str, req: BackgroundRequest):
    slide = slide_map[slide_id]
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = req.color  # e.g. RGBColor(255, 0, 0)
    return {"status": "ok"}

@app.post("/save")
def save_presentation(filename: str = "output.pptx"):
    prs.save(filename)
    return {"status": "saved", "filename": filename}

@app.post("/theme")
def set_theme(theme_name: str):
    # Implementation of setting a theme
    return {"status": "theme set"}

if __name__ == "__main__":
    main()
